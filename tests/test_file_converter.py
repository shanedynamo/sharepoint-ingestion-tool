"""Tests for FileConverter."""

import io
import sys
from unittest.mock import patch, MagicMock

import pytest

sys.path.insert(0, "src")

from utils.file_converter import (
    FileConverter,
    _normalise_ext,
    _safe_filename,
    _extract_pptx_text,
    _extract_xlsx_text,
    SUPPORTED_EXTENSIONS,
)


# ===================================================================
# needs_conversion
# ===================================================================

class TestNeedsConversion:
    @pytest.mark.parametrize("ext", [".pptx", ".ppt", ".xlsx", ".xls"])
    def test_returns_true_for_convertible(self, ext):
        assert FileConverter.needs_conversion(ext) is True

    @pytest.mark.parametrize("ext", [".pdf", ".docx", ".doc", ".txt"])
    def test_returns_false_for_non_convertible(self, ext):
        assert FileConverter.needs_conversion(ext) is False

    def test_returns_false_for_unsupported(self):
        assert FileConverter.needs_conversion(".jpg") is False

    def test_handles_uppercase(self):
        assert FileConverter.needs_conversion(".PPTX") is True

    def test_handles_no_dot_prefix(self):
        assert FileConverter.needs_conversion("xlsx") is True

    def test_handles_whitespace(self):
        assert FileConverter.needs_conversion("  .pptx  ") is True


# ===================================================================
# get_extraction_strategy
# ===================================================================

class TestGetExtractionStrategy:
    @pytest.mark.parametrize("ext,expected", [
        (".pdf", "textract-direct"),
        (".docx", "textract-direct"),
        (".doc", "textract-direct"),
    ])
    def test_textract_direct(self, ext, expected):
        assert FileConverter.get_extraction_strategy(ext) == expected

    @pytest.mark.parametrize("ext", [".pptx", ".ppt", ".xlsx", ".xls"])
    def test_convert_then_textract(self, ext):
        assert FileConverter.get_extraction_strategy(ext) == "convert-then-textract"

    def test_plain_text(self):
        assert FileConverter.get_extraction_strategy(".txt") == "plain-text"

    def test_unsupported_type(self):
        assert FileConverter.get_extraction_strategy(".jpg") == "unsupported"
        assert FileConverter.get_extraction_strategy(".mp4") == "unsupported"

    def test_normalises_input(self):
        assert FileConverter.get_extraction_strategy("PDF") == "textract-direct"
        assert FileConverter.get_extraction_strategy("  XLSX  ") == "convert-then-textract"


# ===================================================================
# convert_to_pdf (LibreOffice)
# ===================================================================

class TestConvertToPdf:
    def test_rejects_non_convertible_type(self):
        with pytest.raises(ValueError, match="does not need conversion"):
            FileConverter.convert_to_pdf(b"data", "file.pdf", ".pdf")

    def test_rejects_unsupported_type(self):
        with pytest.raises(ValueError, match="does not need conversion"):
            FileConverter.convert_to_pdf(b"data", "file.jpg", ".jpg")

    @patch("utils.file_converter.subprocess.run")
    def test_calls_libreoffice(self, mock_run, tmp_path):
        """Verify the subprocess call structure and PDF read-back."""
        mock_result = MagicMock()
        mock_result.returncode = 0

        def side_effect(cmd, **kwargs):
            # Simulate LibreOffice creating a PDF
            outdir = cmd[cmd.index("--outdir") + 1]
            input_file = cmd[-1]
            from pathlib import Path
            pdf_path = Path(outdir) / (Path(input_file).stem + ".pdf")
            pdf_path.write_bytes(b"%PDF-1.4 fake")
            return mock_result

        mock_run.side_effect = side_effect

        result = FileConverter.convert_to_pdf(
            b"pptx-content", "slides.pptx", ".pptx",
        )

        assert result == b"%PDF-1.4 fake"
        mock_run.assert_called_once()
        call_args = mock_run.call_args
        cmd = call_args[0][0]
        assert cmd[0] == "libreoffice"
        assert "--headless" in cmd
        assert "--convert-to" in cmd
        assert cmd[cmd.index("--convert-to") + 1] == "pdf"

    @patch("utils.file_converter.subprocess.run")
    def test_raises_on_nonzero_exit(self, mock_run):
        mock_run.return_value = MagicMock(
            returncode=1, stderr="conversion error",
        )
        with pytest.raises(RuntimeError, match="PDF conversion failed"):
            FileConverter.convert_to_pdf(b"data", "bad.xlsx", ".xlsx")

    @patch("utils.file_converter.subprocess.run")
    def test_raises_when_pdf_not_created(self, mock_run):
        mock_run.return_value = MagicMock(returncode=0, stderr="")
        # subprocess returns 0 but no PDF created
        with pytest.raises(RuntimeError, match="Expected PDF output not found"):
            FileConverter.convert_to_pdf(b"data", "bad.pptx", ".pptx")


# ===================================================================
# convert_to_pdf_lambda
# ===================================================================

class TestConvertToPdfLambda:
    def test_rejects_legacy_ppt(self):
        with pytest.raises(ValueError, match="Legacy format"):
            FileConverter.convert_to_pdf_lambda(b"data", "old.ppt", ".ppt")

    def test_rejects_legacy_xls(self):
        with pytest.raises(ValueError, match="Legacy format"):
            FileConverter.convert_to_pdf_lambda(b"data", "old.xls", ".xls")

    def test_rejects_unsupported_type(self):
        with pytest.raises(ValueError, match="not supported for Lambda"):
            FileConverter.convert_to_pdf_lambda(b"data", "img.jpg", ".jpg")

    def test_extracts_pptx_text(self):
        """Create a real PPTX in-memory and extract text."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # title+content
        slide.shapes.title.text = "Slide Title"
        slide.placeholders[1].text = "Bullet point text"

        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Second Slide"
        slide2.placeholders[1].text = "More content"

        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        result = FileConverter.convert_to_pdf_lambda(
            pptx_bytes, "test.pptx", ".pptx",
        )

        text = result.decode("utf-8")
        assert "Slide Title" in text
        assert "Bullet point text" in text
        assert "Second Slide" in text
        assert "More content" in text
        assert "--- Slide 1 ---" in text
        assert "--- Slide 2 ---" in text

    def test_extracts_xlsx_text(self):
        """Create a real XLSX in-memory and extract text."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["A2"] = "Alice"
        ws["B2"] = 42
        ws["A3"] = "Bob"
        ws["B3"] = 99

        ws2 = wb.create_sheet("Summary")
        ws2["A1"] = "Total"
        ws2["B1"] = 141

        buf = io.BytesIO()
        wb.save(buf)
        xlsx_bytes = buf.getvalue()

        result = FileConverter.convert_to_pdf_lambda(
            xlsx_bytes, "test.xlsx", ".xlsx",
        )

        text = result.decode("utf-8")
        assert "--- Data ---" in text
        assert "Name\tValue" in text
        assert "Alice\t42" in text
        assert "Bob\t99" in text
        assert "--- Summary ---" in text
        assert "Total\t141" in text

    def test_pptx_empty_presentation(self):
        """Empty PPTX returns empty bytes."""
        from pptx import Presentation

        prs = Presentation()
        buf = io.BytesIO()
        prs.save(buf)

        result = FileConverter.convert_to_pdf_lambda(
            buf.getvalue(), "empty.pptx", ".pptx",
        )
        assert result == b""

    def test_xlsx_empty_workbook(self):
        """Empty XLSX returns empty bytes."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        # Leave all cells empty
        buf = io.BytesIO()
        wb.save(buf)

        result = FileConverter.convert_to_pdf_lambda(
            buf.getvalue(), "empty.xlsx", ".xlsx",
        )
        assert result == b""

    def test_pptx_with_table(self):
        """PPTX with a table shape extracts cell text."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

        rows, cols = 2, 2
        table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(4), Inches(2)).table
        table.cell(0, 0).text = "Header1"
        table.cell(0, 1).text = "Header2"
        table.cell(1, 0).text = "Val1"
        table.cell(1, 1).text = "Val2"

        buf = io.BytesIO()
        prs.save(buf)

        result = FileConverter.convert_to_pdf_lambda(
            buf.getvalue(), "tables.pptx", ".pptx",
        )
        text = result.decode("utf-8")
        assert "Header1\tHeader2" in text
        assert "Val1\tVal2" in text


# ===================================================================
# _normalise_ext
# ===================================================================

class TestNormaliseExt:
    def test_already_normalised(self):
        assert _normalise_ext(".pdf") == ".pdf"

    def test_adds_dot(self):
        assert _normalise_ext("pdf") == ".pdf"

    def test_lowercases(self):
        assert _normalise_ext(".PPTX") == ".pptx"

    def test_strips_whitespace(self):
        assert _normalise_ext("  .xlsx  ") == ".xlsx"


# ===================================================================
# _safe_filename
# ===================================================================

class TestSafeFilename:
    def test_normal_filename(self):
        assert _safe_filename("report.pptx", ".pptx") == "report.pptx"

    def test_special_characters_removed(self):
        result = _safe_filename("my<>file|name.xlsx", ".xlsx")
        assert "<" not in result
        assert ">" not in result
        assert "|" not in result
        assert result.endswith(".xlsx")

    def test_preserves_hyphens_underscores_spaces(self):
        result = _safe_filename("my-file_name v2.pptx", ".pptx")
        assert result == "my-file_name v2.pptx"

    def test_empty_stem_uses_default(self):
        result = _safe_filename(".pptx", ".pptx")
        assert result == "document.pptx"


# ===================================================================
# SUPPORTED_EXTENSIONS constant
# ===================================================================

class TestSupportedExtensions:
    def test_contains_all_types(self):
        expected = {
            ".pptx", ".ppt", ".xlsx", ".xls",
            ".pdf", ".docx", ".doc",
            ".txt",
        }
        assert SUPPORTED_EXTENSIONS == expected
