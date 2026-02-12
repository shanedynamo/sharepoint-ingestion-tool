"""Generate test fixture files for integration tests.

Creates sample documents (PDF, DOCX, PPTX, XLSX) in tests/fixtures/.
Can be run standalone or imported by conftest.py fixtures.

Dependencies:
    fpdf2       - PDF generation
    python-docx - DOCX generation
    python-pptx - PPTX generation (project dependency)
    openpyxl    - XLSX generation (project dependency)
"""

import os
from pathlib import Path

FIXTURES_DIR = Path(__file__).parent


def generate_pdf(output_path: Path | None = None) -> Path:
    """Generate a 2-page PDF with sample text."""
    path = output_path or FIXTURES_DIR / "sample.pdf"
    if path.exists():
        return path

    try:
        from fpdf import FPDF

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)

        # Page 1
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 16)
        pdf.cell(0, 10, "SharePoint Integration Test Document", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 12)
        pdf.ln(5)
        pdf.multi_cell(0, 7, (
            "This is a sample PDF document used for integration testing of the "
            "SharePoint ingestion pipeline. It contains two pages of text content "
            "that will be processed by the extraction pipeline."
        ))
        pdf.ln(5)
        pdf.multi_cell(0, 7, (
            "The document registry should track this file with textract_status "
            "set to 'pending' initially, then updated as processing completes."
        ))

        # Page 2
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 14)
        pdf.cell(0, 10, "Page Two - Additional Content", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 12)
        pdf.ln(5)
        pdf.multi_cell(0, 7, (
            "This second page contains additional content to verify that "
            "multi-page document extraction works correctly."
        ))

        pdf.output(str(path))
    except ImportError:
        # Fallback: craft a minimal valid PDF by hand
        _write_minimal_pdf(path)

    return path


def generate_docx(output_path: Path | None = None) -> Path:
    """Generate a Word document with sample content."""
    path = output_path or FIXTURES_DIR / "sample.docx"
    if path.exists():
        return path

    try:
        from docx import Document

        doc = Document()
        doc.add_heading("SharePoint Integration Test", level=1)
        doc.add_paragraph(
            "This is a sample Word document for integration testing. "
            "It contains headings and paragraphs to verify DOCX processing."
        )
        doc.add_heading("Section Two", level=2)
        doc.add_paragraph(
            "Additional content in a second section to test multi-section extraction."
        )
        doc.save(str(path))
    except ImportError:
        # Fallback: create a minimal DOCX (it's a ZIP file with XML)
        _write_minimal_docx(path)

    return path


def generate_pptx(output_path: Path | None = None) -> Path:
    """Generate a 3-slide PPTX presentation."""
    path = output_path or FIXTURES_DIR / "sample.pptx"
    if path.exists():
        return path

    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Integration Test Presentation"
    slide1.placeholders[1].text = "Generated for LocalStack testing"

    # Slide 2: Content with bullet points
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Test Content"
    body = slide2.placeholders[1]
    body.text = "First bullet point"
    body.text_frame.add_paragraph().text = "Second bullet point"
    body.text_frame.add_paragraph().text = "Third bullet point"

    # Slide 3: More content
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Additional Data"
    body3 = slide3.placeholders[1]
    body3.text = "This slide contains additional test data"

    prs.save(str(path))
    return path


def generate_xlsx(output_path: Path | None = None) -> Path:
    """Generate a 2-sheet XLSX spreadsheet."""
    path = output_path or FIXTURES_DIR / "sample.xlsx"
    if path.exists():
        return path

    from openpyxl import Workbook

    wb = Workbook()

    # Sheet 1: Employees
    ws1 = wb.active
    ws1.title = "Employees"
    ws1.append(["Name", "Department", "Start Date"])
    ws1.append(["Alice Johnson", "Engineering", "2024-01-15"])
    ws1.append(["Bob Smith", "Marketing", "2023-06-01"])
    ws1.append(["Carol Williams", "HR", "2024-03-20"])

    # Sheet 2: Projects
    ws2 = wb.create_sheet("Projects")
    ws2.append(["Project", "Status", "Budget"])
    ws2.append(["SharePoint Migration", "Active", "50000"])
    ws2.append(["AI Integration", "Planning", "120000"])

    wb.save(str(path))
    return path


def generate_all() -> dict[str, Path]:
    """Generate all fixture files and return a mapping of name -> path."""
    return {
        "pdf": generate_pdf(),
        "docx": generate_docx(),
        "pptx": generate_pptx(),
        "xlsx": generate_xlsx(),
    }


# -----------------------------------------------------------------------
# Fallback generators (no external dependencies)
# -----------------------------------------------------------------------

def _write_minimal_pdf(path: Path) -> None:
    """Write a minimal valid PDF file without any dependencies."""
    content = (
        "%PDF-1.4\n"
        "1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
        "2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
        "3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        "/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>\nendobj\n"
        "4 0 obj\n<< /Length 44 >>\nstream\n"
        "BT /F1 12 Tf 100 700 Td (Test PDF) Tj ET\n"
        "endstream\nendobj\n"
        "5 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n"
        "xref\n0 6\n"
        "0000000000 65535 f \n"
        "0000000009 00000 n \n"
        "0000000058 00000 n \n"
        "0000000115 00000 n \n"
        "0000000266 00000 n \n"
        "0000000360 00000 n \n"
        "trailer\n<< /Size 6 /Root 1 0 R >>\nstartxref\n436\n%%EOF\n"
    )
    path.write_text(content)


def _write_minimal_docx(path: Path) -> None:
    """Write a minimal valid DOCX (Office Open XML) without python-docx."""
    import zipfile

    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        '</Types>'
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
        'Target="word/document.xml"/>'
        '</Relationships>'
    )
    document = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        '<w:body><w:p><w:r><w:t>Test DOCX content</w:t></w:r></w:p></w:body>'
        '</w:document>'
    )

    with zipfile.ZipFile(str(path), "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types)
        zf.writestr("_rels/.rels", rels)
        zf.writestr("word/document.xml", document)


if __name__ == "__main__":
    paths = generate_all()
    for name, p in paths.items():
        print(f"  {name}: {p}")
    print("Done.")
