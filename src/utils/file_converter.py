"""Convert Office formats to PDF for Textract processing.

Provides two conversion paths:

* **EC2 / Docker** – ``convert_to_pdf`` shells out to LibreOffice headless.
* **Lambda fallback** – ``convert_to_pdf_lambda`` uses *python-pptx* and
  *openpyxl* to extract text directly (no native binary required).

A strategy helper (``get_extraction_strategy``) maps each file type to the
correct downstream pipeline step.
"""

import io
import logging
import subprocess
import tempfile
from pathlib import Path

logger = logging.getLogger(__name__)

# --------------------------------------------------------------------- #
# Extension sets
# --------------------------------------------------------------------- #

# Types that need conversion before Textract can process them.
_NEEDS_CONVERSION = {".pptx", ".ppt", ".xlsx", ".xls"}

# Types that Textract handles natively (no conversion needed).
_TEXTRACT_DIRECT = {".pdf", ".docx", ".doc"}

# Plain-text types – no Textract needed at all.
_PLAIN_TEXT = {".txt"}

# All supported types (union).
SUPPORTED_EXTENSIONS = _NEEDS_CONVERSION | _TEXTRACT_DIRECT | _PLAIN_TEXT


class FileConverter:
    """Converts Office documents and classifies extraction strategies."""

    # ----------------------------------------------------------------- #
    # Classification helpers
    # ----------------------------------------------------------------- #

    @staticmethod
    def needs_conversion(file_type: str) -> bool:
        """Return *True* if *file_type* must be converted before Textract.

        Normalises the extension to lowercase with a leading dot.
        """
        ext = _normalise_ext(file_type)
        return ext in _NEEDS_CONVERSION

    @staticmethod
    def get_extraction_strategy(file_type: str) -> str:
        """Return the extraction strategy name for *file_type*.

        Returns one of:

        * ``"textract-direct"`` – PDF / DOCX handled by Textract as-is.
        * ``"convert-then-textract"`` – Convert to PDF first (EC2 path).
        * ``"direct-extract"`` – Extract text in-process (Lambda path).
        * ``"plain-text"`` – Read the bytes as UTF-8 text.
        * ``"unsupported"`` – File type is not handled.
        """
        ext = _normalise_ext(file_type)
        if ext in _TEXTRACT_DIRECT:
            return "textract-direct"
        if ext in _NEEDS_CONVERSION:
            return "convert-then-textract"
        if ext in _PLAIN_TEXT:
            return "plain-text"
        return "unsupported"

    # ----------------------------------------------------------------- #
    # EC2 conversion (LibreOffice headless)
    # ----------------------------------------------------------------- #

    @staticmethod
    def convert_to_pdf(
        content: bytes,
        filename: str,
        file_type: str,
    ) -> bytes:
        """Convert an Office document to PDF via LibreOffice headless.

        Requires ``libreoffice`` on ``$PATH`` (available in the EC2 / Docker
        image).  Raises :class:`RuntimeError` on failure.
        """
        ext = _normalise_ext(file_type)
        if ext not in _NEEDS_CONVERSION:
            raise ValueError(
                f"File type '{ext}' does not need conversion. "
                f"Convertible types: {sorted(_NEEDS_CONVERSION)}"
            )

        with tempfile.TemporaryDirectory() as tmpdir:
            # Write source file
            safe_name = _safe_filename(filename, ext)
            input_path = Path(tmpdir) / safe_name
            input_path.write_bytes(content)

            result = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", tmpdir,
                    str(input_path),
                ],
                capture_output=True,
                text=True,
                timeout=120,
            )

            if result.returncode != 0:
                logger.error(
                    "LibreOffice conversion failed for %s: %s",
                    filename, result.stderr,
                )
                raise RuntimeError(
                    f"PDF conversion failed for {filename}: {result.stderr}"
                )

            pdf_path = input_path.with_suffix(".pdf")
            if not pdf_path.exists():
                raise RuntimeError(
                    f"Expected PDF output not found: {pdf_path}"
                )

            return pdf_path.read_bytes()

    # ----------------------------------------------------------------- #
    # Lambda fallback (pure-Python extraction)
    # ----------------------------------------------------------------- #

    @staticmethod
    def convert_to_pdf_lambda(
        content: bytes,
        filename: str,
        file_type: str,
    ) -> bytes:
        """Extract text from Office files using pure-Python libraries.

        Returns UTF-8 encoded text (not a real PDF) suitable for direct
        indexing.  This is the Lambda fallback when LibreOffice is not
        available.

        Supported types: ``.pptx``, ``.xlsx``.  For ``.ppt`` / ``.xls``
        (legacy binary), raises :class:`ValueError` since pure-Python
        libraries cannot read them.
        """
        ext = _normalise_ext(file_type)

        if ext == ".pptx":
            return _extract_pptx_text(content, filename)
        if ext == ".xlsx":
            return _extract_xlsx_text(content, filename)
        if ext in (".ppt", ".xls"):
            raise ValueError(
                f"Legacy format '{ext}' is not supported in Lambda mode. "
                "Use convert_to_pdf (LibreOffice) on EC2 instead."
            )
        raise ValueError(
            f"File type '{ext}' is not supported for Lambda conversion. "
            f"Supported: .pptx, .xlsx"
        )


# --------------------------------------------------------------------- #
# Internal helpers
# --------------------------------------------------------------------- #

def _normalise_ext(file_type: str) -> str:
    """Ensure the extension is lowercase and starts with a dot."""
    ext = file_type.strip().lower()
    if not ext.startswith("."):
        ext = f".{ext}"
    return ext


def _safe_filename(filename: str, ext: str) -> str:
    """Return a filesystem-safe filename with the correct extension."""
    stem = Path(filename).stem
    # Path(".pptx").stem == ".pptx" — treat as empty
    if not stem or stem.startswith("."):
        stem = "document"
    # Remove characters problematic on most filesystems
    safe = "".join(c for c in stem if c.isalnum() or c in "-_ ")
    return f"{safe or 'document'}{ext}"


def _extract_pptx_text(content: bytes, filename: str) -> bytes:
    """Extract all text from a PPTX file using python-pptx."""
    from pptx import Presentation  # lazy import

    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = "\t".join(
                        cell.text.strip() for cell in row.cells
                    )
                    if row_text.strip():
                        slide_texts.append(row_text)

        if slide_texts:
            parts.append(f"--- Slide {slide_idx} ---")
            parts.extend(slide_texts)

    text = "\n".join(parts)
    logger.info("Extracted %d characters from PPTX: %s", len(text), filename)
    return text.encode("utf-8")


def _extract_xlsx_text(content: bytes, filename: str) -> bytes:
    """Extract all cell values from an XLSX file using openpyxl."""
    from openpyxl import load_workbook  # lazy import

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts: list[str] = []

    for sheet in wb.worksheets:
        sheet_texts: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(v).strip() for v in row if v is not None]
            if cells:
                sheet_texts.append("\t".join(cells))

        if sheet_texts:
            parts.append(f"--- {sheet.title} ---")
            parts.extend(sheet_texts)

    wb.close()
    text = "\n".join(parts)
    logger.info("Extracted %d characters from XLSX: %s", len(text), filename)
    return text.encode("utf-8")
